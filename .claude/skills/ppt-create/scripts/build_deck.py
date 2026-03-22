#!/usr/bin/env python3
"""Build a PowerPoint deck from a JSON manifest.

Usage:
    cd .claude/skills/ppt-create/references && uv run python ../scripts/build_deck.py --manifest <path> --output <path>
"""

import argparse
import inspect
import json
import logging
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")

# ---------------------------------------------------------------------------
# Palette registry
# ---------------------------------------------------------------------------

PALETTES = {
    "midnight": {"bg": "1E1E2E", "fg": "CDD6F4", "accent": "89B4FA", "accent2": "A6E3A1", "muted": "6C7086", "highlight": "F9E2AF", "error": "F38BA8"},
    "charcoal": {"bg": "36454F", "fg": "F2F2F2", "accent": "4FC3F7", "accent2": "81C784", "muted": "90A4AE", "highlight": "FFD54F", "error": "EF5350"},
    "ocean": {"bg": "0D1B2A", "fg": "E0E1DD", "accent": "1B9AAA", "accent2": "06D6A0", "muted": "415A77", "highlight": "FFD166", "error": "EF476F"},
    "sage-calm": {"bg": "2D3A2D", "fg": "E8EDE8", "accent": "8FBC8F", "accent2": "B8D4A3", "muted": "6B7F6B", "highlight": "E6D99E", "error": "D4736D"},
    "berry-cream": {"bg": "2C1E2E", "fg": "F4E8F0", "accent": "C084D8", "accent2": "E8A0BF", "muted": "7A5C82", "highlight": "F2D58A", "error": "E85A71"},
    "cherry-bold": {"bg": "2B1215", "fg": "F5E6E8", "accent": "E74C3C", "accent2": "F39C12", "muted": "7B4A50", "highlight": "F7DC6F", "error": "FF6B6B"},
    "forest-moss": {"bg": "1A2E1A", "fg": "E0EBD8", "accent": "4CAF50", "accent2": "8BC34A", "muted": "5D7A5D", "highlight": "CDDC39", "error": "EF5350"},
    "coral-energy": {"bg": "2E1F1A", "fg": "FFF0E8", "accent": "FF7043", "accent2": "FFB74D", "muted": "8D6E63", "highlight": "FFE082", "error": "E53935"},
    "ocean-gradient": {"bg": "0A2540", "fg": "E6F0FA", "accent": "00BCD4", "accent2": "26C6DA", "muted": "456880", "highlight": "FFD740", "error": "FF5252"},
    "midnight-exec": {"bg": "1A1A2E", "fg": "EAEAF2", "accent": "6C63FF", "accent2": "48C9B0", "muted": "5A5A82", "highlight": "F0C040", "error": "E84A5F"},
    "warm-terracotta": {"bg": "3E2723", "fg": "EFEBE9", "accent": "D4836A", "accent2": "A1887F", "muted": "795548", "highlight": "FFE0B2", "error": "E57373"},
    "teal-trust": {"bg": "0D2B2B", "fg": "E0F2F1", "accent": "009688", "accent2": "4DB6AC", "muted": "4A7A78", "highlight": "FFD54F", "error": "EF5350"},
}


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb(hex_str):
    """Convert a 6-char hex string to RGBColor."""
    try:
        return RGBColor.from_string(hex_str)
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def _set_bg(slide, palette):
    """Set slide background to palette bg colour."""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(palette["bg"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def _add_textbox(slide, left, top, width, height, text, font_name, font_size,
                 bold=False, italic=False, color_hex=None, alignment=None, word_wrap=True):
    """Add a styled text box to a slide."""
    try:
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = font_name
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        if color_hex:
            p.font.color.rgb = _rgb(color_hex)
        if alignment:
            p.alignment = alignment
        return txBox
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def _add_bullets(slide, left, top, width, height, items, font_name, font_size,
                 color_hex=None, alignment=None):
    """Add a bulleted list text box."""
    try:
        txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.name = font_name
            p.font.size = Pt(font_size)
            p.level = 0
            if color_hex:
                p.font.color.rgb = _rgb(color_hex)
            if alignment:
                p.alignment = alignment
            p.space_after = Pt(4)
        return txBox
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def _resolve_image(image_path, manifest_dir):
    """Resolve image path relative to manifest directory if not absolute."""
    try:
        if os.path.isabs(image_path):
            return image_path
        return os.path.normpath(os.path.join(manifest_dir, image_path))
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# Slide renderers
# ---------------------------------------------------------------------------

def render_title(prs, slide_data, palette, manifest_dir):
    """Render a title slide."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        # Optional background image
        if slide_data.get("image_path"):
            img = _resolve_image(slide_data["image_path"], manifest_dir)
            if os.path.exists(img):
                slide.shapes.add_picture(img, Inches(0), Inches(0), width=Inches(10))

        _add_textbox(slide, 0.8, 2.0, 8.4, 1.5, slide_data["title"],
                     "Georgia", 44, bold=True, color_hex=palette["fg"])

        if slide_data.get("subtitle"):
            _add_textbox(slide, 0.8, 3.6, 8.4, 0.8, slide_data["subtitle"],
                         "Calibri", 20, color_hex=palette["fg"])

        if slide_data.get("date"):
            _add_textbox(slide, 0.8, 6.5, 8.4, 0.5, slide_data["date"],
                         "Calibri", 12, italic=True, color_hex=palette["muted"])

        logger.info("Added title slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_section_header(prs, slide_data, palette, manifest_dir):
    """Render a section header slide."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        _add_textbox(slide, 0.8, 2.8, 8.4, 1.2, slide_data["title"],
                     "Georgia", 36, bold=True, color_hex=palette["fg"],
                     alignment=PP_ALIGN.CENTER)

        if slide_data.get("subtitle"):
            _add_textbox(slide, 1.5, 4.2, 7.0, 0.6, slide_data["subtitle"],
                         "Calibri", 16, color_hex=palette["muted"],
                         alignment=PP_ALIGN.CENTER)

        logger.info("Added section_header slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_content(prs, slide_data, palette, manifest_dir):
    """Render a content slide with body text and optional bullets."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        _add_textbox(slide, 0.8, 0.4, 8.4, 0.8, slide_data["title"],
                     "Georgia", 36, bold=True, color_hex=palette["fg"])

        _add_textbox(slide, 0.8, 1.5, 8.4, 2.5, slide_data["body"],
                     "Calibri", 16, color_hex=palette["fg"])

        if slide_data.get("bullets"):
            _add_bullets(slide, 1.0, 4.2, 8.0, 2.8, slide_data["bullets"],
                         "Calibri", 14, color_hex=palette["fg"])

        logger.info("Added content slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_diagram(prs, slide_data, palette, manifest_dir):
    """Render a diagram slide with image and optional caption."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        _add_textbox(slide, 0.8, 0.4, 8.4, 0.8, slide_data["title"],
                     "Georgia", 28, bold=True, color_hex=palette["fg"])

        img = _resolve_image(slide_data["image_path"], manifest_dir)
        slide.shapes.add_picture(img, left=Inches(0.75), top=Inches(1.5), width=Inches(8.5))

        if slide_data.get("caption"):
            _add_textbox(slide, 0.8, 6.8, 8.4, 0.4, slide_data["caption"],
                         "Calibri", 10, italic=True, color_hex=palette["muted"],
                         alignment=PP_ALIGN.CENTER)

        logger.info("Added diagram slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_diagram_full(prs, slide_data, palette, manifest_dir):
    """Render a full-bleed diagram slide."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        img = _resolve_image(slide_data["image_path"], manifest_dir)
        slide.shapes.add_picture(img, left=Inches(0), top=Inches(0), width=Inches(10))

        if slide_data.get("title"):
            _add_textbox(slide, 0.5, 0.3, 5.0, 0.6, slide_data["title"],
                         "Georgia", 20, bold=True, color_hex=palette["fg"])

        logger.info("Added diagram_full slide")
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_metrics(prs, slide_data, palette, manifest_dir):
    """Render a metrics slide with large stat callouts."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        _add_textbox(slide, 0.8, 0.4, 8.4, 0.8, slide_data["title"],
                     "Georgia", 36, bold=True, color_hex=palette["fg"])

        metrics = slide_data["metrics"]
        n = len(metrics)
        col_width = 8.4 / n if n > 0 else 8.4

        for i, m in enumerate(metrics):
            x = 0.8 + i * col_width
            # Value
            _add_textbox(slide, x, 2.0, col_width, 1.5, str(m["value"]),
                         "Georgia", 60, bold=True, color_hex=palette["accent"],
                         alignment=PP_ALIGN.CENTER)
            # Label
            _add_textbox(slide, x, 3.8, col_width, 0.8, m["label"],
                         "Calibri", 14, color_hex=palette["muted"],
                         alignment=PP_ALIGN.CENTER)

        logger.info("Added metrics slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_two_column(prs, slide_data, palette, manifest_dir):
    """Render a two-column slide."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        _add_textbox(slide, 0.8, 0.4, 8.4, 0.8, slide_data["title"],
                     "Georgia", 36, bold=True, color_hex=palette["fg"])

        left_top = 1.5
        right_top = 1.5

        # Left image
        if slide_data.get("left_image"):
            img = _resolve_image(slide_data["left_image"], manifest_dir)
            if os.path.exists(img):
                slide.shapes.add_picture(img, left=Inches(0.5), top=Inches(1.5), width=Inches(4.2))
                left_top = 4.0

        # Right image
        if slide_data.get("right_image"):
            img = _resolve_image(slide_data["right_image"], manifest_dir)
            if os.path.exists(img):
                slide.shapes.add_picture(img, left=Inches(5.3), top=Inches(1.5), width=Inches(4.2))
                right_top = 4.0

        _add_textbox(slide, 0.5, left_top, 4.2, 5.0 - (left_top - 1.5),
                     slide_data["left_body"], "Calibri", 14, color_hex=palette["fg"])

        _add_textbox(slide, 5.3, right_top, 4.2, 5.0 - (right_top - 1.5),
                     slide_data["right_body"], "Calibri", 14, color_hex=palette["fg"])

        logger.info("Added two_column slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_comparison(prs, slide_data, palette, manifest_dir):
    """Render a comparison (before/after, pros/cons) slide."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        _add_textbox(slide, 0.8, 0.4, 8.4, 0.8, slide_data["title"],
                     "Georgia", 36, bold=True, color_hex=palette["fg"])

        # Left header + items
        _add_textbox(slide, 0.5, 1.5, 4.2, 0.5, slide_data["left_label"],
                     "Georgia", 20, bold=True, color_hex=palette["accent"])
        _add_bullets(slide, 0.7, 2.2, 4.0, 3.0, slide_data["left_items"],
                     "Calibri", 14, color_hex=palette["fg"])

        # Right header + items
        _add_textbox(slide, 5.3, 1.5, 4.2, 0.5, slide_data["right_label"],
                     "Georgia", 20, bold=True, color_hex=palette["accent2"])
        _add_bullets(slide, 5.5, 2.2, 4.0, 3.0, slide_data["right_items"],
                     "Calibri", 14, color_hex=palette["fg"])

        # Optional images
        if slide_data.get("left_image"):
            img = _resolve_image(slide_data["left_image"], manifest_dir)
            if os.path.exists(img):
                slide.shapes.add_picture(img, left=Inches(0.5), top=Inches(5.5), width=Inches(4.2))
        if slide_data.get("right_image"):
            img = _resolve_image(slide_data["right_image"], manifest_dir)
            if os.path.exists(img):
                slide.shapes.add_picture(img, left=Inches(5.3), top=Inches(5.5), width=Inches(4.2))

        logger.info("Added comparison slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def render_closing(prs, slide_data, palette, manifest_dir):
    """Render a closing slide."""
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide, palette)

        _add_textbox(slide, 0.8, 2.0, 8.4, 1.2, slide_data["title"],
                     "Georgia", 40, bold=True, color_hex=palette["fg"],
                     alignment=PP_ALIGN.CENTER)

        if slide_data.get("bullets"):
            _add_bullets(slide, 1.5, 3.5, 7.0, 2.5, slide_data["bullets"],
                         "Calibri", 16, color_hex=palette["fg"],
                         alignment=PP_ALIGN.CENTER)

        if slide_data.get("call_to_action"):
            _add_textbox(slide, 1.5, 6.2, 7.0, 0.6, slide_data["call_to_action"],
                         "Calibri", 18, bold=True, color_hex=palette["accent"],
                         alignment=PP_ALIGN.CENTER)

        logger.info("Added closing slide: %s", slide_data["title"])
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# Dispatcher
# ---------------------------------------------------------------------------

RENDERERS = {
    "title": render_title,
    "section_header": render_section_header,
    "content": render_content,
    "diagram": render_diagram,
    "diagram_full": render_diagram_full,
    "metrics": render_metrics,
    "two_column": render_two_column,
    "comparison": render_comparison,
    "closing": render_closing,
}


def build_deck(manifest_path, output_path):
    """Build a PowerPoint deck from a JSON manifest file."""
    try:
        manifest_path = os.path.abspath(manifest_path)
        output_path = os.path.abspath(output_path)
        manifest_dir = os.path.dirname(manifest_path)

        with open(manifest_path, "r", encoding="utf-8") as f:
            manifest = json.load(f)

        palette_name = manifest.get("palette", "midnight")
        palette = PALETTES.get(palette_name)
        if palette is None:
            raise ValueError(f"Unknown palette: {palette_name}. Available: {list(PALETTES.keys())}")

        prs = Presentation()
        # Standard 10" x 7.5" slide dimensions (default)

        slides = manifest.get("slides", [])
        if not slides:
            raise ValueError("Manifest contains no slides")

        for i, slide_data in enumerate(slides):
            slide_type = slide_data.get("type")
            if slide_type not in RENDERERS:
                raise ValueError(f"Slide {i}: unknown type '{slide_type}'. Available: {list(RENDERERS.keys())}")
            RENDERERS[slide_type](prs, slide_data, palette, manifest_dir)

        prs.save(output_path)
        logger.info("Saved deck to %s (%d slides)", output_path, len(slides))
        return output_path
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main():
    """CLI entry point."""
    try:
        parser = argparse.ArgumentParser(description="Build a PowerPoint deck from a JSON manifest.")
        parser.add_argument("--manifest", required=True, help="Path to JSON manifest file")
        parser.add_argument("--output", required=True, help="Output .pptx file path")
        args = parser.parse_args()

        build_deck(args.manifest, args.output)
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


if __name__ == "__main__":
    main()
