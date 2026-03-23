#!/usr/bin/env python3
"""End-to-end integration test for the ppt-create pipeline.

Tests the full flow: markdown -> parse_report -> manifest -> build_deck -> pptx -> validate_deck.

Usage:
    cd .claude/skills/ppt-create/references && uv run python ../scripts/test_e2e.py
"""

import inspect
import json
import logging
import os
import shutil
import sys
import tempfile
from typing import Any

# Ensure scripts directory is importable
scripts_dir = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, scripts_dir)

from parse_report import parse_markdown
from build_deck import build_deck, PALETTES
from validate_deck import validate_deck

from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.dml.color import RGBColor

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")


# ---------------------------------------------------------------------------
# Test helpers
# ---------------------------------------------------------------------------

def create_test_image(path: str, width: int = 200, height: int = 100) -> str:
    """Create a solid-colour PNG test image.

    Args:
        path: File path for the output PNG.
        width: Image width in pixels.
        height: Image height in pixels.

    Returns:
        Absolute path to the created image.
    """
    try:
        # Create a minimal valid PNG without Pillow
        # Using a simple approach: write raw PNG bytes for a solid colour rectangle
        import struct
        import zlib

        def _create_png(w: int, h: int, r: int, g: int, b: int) -> bytes:
            """Generate a minimal PNG file as bytes."""
            # PNG signature
            sig = b'\x89PNG\r\n\x1a\n'

            # IHDR chunk
            ihdr_data = struct.pack('>IIBBBBB', w, h, 8, 2, 0, 0, 0)  # 8-bit RGB
            ihdr_crc = zlib.crc32(b'IHDR' + ihdr_data) & 0xffffffff
            ihdr = struct.pack('>I', 13) + b'IHDR' + ihdr_data + struct.pack('>I', ihdr_crc)

            # IDAT chunk - raw pixel data
            raw_data = b''
            for _row in range(h):
                raw_data += b'\x00'  # filter byte
                raw_data += bytes([r, g, b]) * w
            compressed = zlib.compress(raw_data)
            idat_crc = zlib.crc32(b'IDAT' + compressed) & 0xffffffff
            idat = struct.pack('>I', len(compressed)) + b'IDAT' + compressed + struct.pack('>I', idat_crc)

            # IEND chunk
            iend_crc = zlib.crc32(b'IEND') & 0xffffffff
            iend = struct.pack('>I', 0) + b'IEND' + struct.pack('>I', iend_crc)

            return sig + ihdr + idat + iend

        png_bytes = _create_png(width, height, 70, 130, 180)  # Steel blue
        abs_path = os.path.abspath(path)
        with open(abs_path, 'wb') as f:
            f.write(png_bytes)
        return abs_path
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def create_test_markdown(image_filename: str) -> str:
    """Create test markdown content with 3+ sections, bullets, and an image reference.

    Args:
        image_filename: Filename of the test image to reference.

    Returns:
        Markdown string.
    """
    try:
        return f"""# Integration Test Report

This is the subtitle for the presentation.

## Key Findings

### Performance Results

Our tests showed significant improvements across all metrics.

- Response time reduced by 40%
- Throughput increased by 60%
- Error rate dropped to 0.1%

### Architecture Diagram

![System architecture]({image_filename})

### Summary

This concludes our integration test report.

- All systems operational
- No critical issues found
"""
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# Test runner
# ---------------------------------------------------------------------------

def run_e2e_test() -> bool:
    """Run the full end-to-end pipeline test.

    Returns:
        True if all assertions pass, False otherwise.
    """
    try:
        test_dir = tempfile.mkdtemp(prefix="ppt_e2e_test_")
        print(f"\n{'='*60}")
        print(f"E2E Test - working directory: {test_dir}")
        print(f"{'='*60}")

        passed = True

        try:
            # ---------------------------------------------------------------
            # Step 1: Create test image
            # ---------------------------------------------------------------
            print("\n--- Step 1: Create test PNG image ---")
            image_path = create_test_image(os.path.join(test_dir, "test_image.png"))
            assert os.path.exists(image_path), "Test image was not created"
            print(f"  Created: {image_path} ({os.path.getsize(image_path)} bytes)")

            # ---------------------------------------------------------------
            # Step 2: Create test markdown
            # ---------------------------------------------------------------
            print("\n--- Step 2: Create test markdown ---")
            md_content = create_test_markdown("test_image.png")
            md_path = os.path.join(test_dir, "test_report.md")
            with open(md_path, "w", encoding="utf-8") as f:
                f.write(md_content)
            print(f"  Created: {md_path}")

            # ---------------------------------------------------------------
            # Step 3: Run parse_report -> manifest JSON
            # ---------------------------------------------------------------
            print("\n--- Step 3: Parse markdown -> manifest ---")
            palette_name = "midnight"
            manifest = parse_markdown(md_content, images_dir=test_dir,
                                      archetype="general", palette=palette_name)
            manifest_path = os.path.join(test_dir, "manifest.json")
            with open(manifest_path, "w", encoding="utf-8") as f:
                json.dump(manifest, f, indent=2)
            print(f"  Manifest slides: {len(manifest['slides'])}")
            print(f"  Palette: {manifest['palette']}")
            for i, s in enumerate(manifest["slides"]):
                print(f"    Slide {i+1}: type={s['type']}, title={s.get('title', 'N/A')}")

            assert len(manifest["slides"]) >= 3, \
                f"Expected at least 3 slides, got {len(manifest['slides'])}"

            # ---------------------------------------------------------------
            # Step 4: Run build_deck -> .pptx
            # ---------------------------------------------------------------
            print("\n--- Step 4: Build deck -> .pptx ---")
            pptx_path = os.path.join(test_dir, "test_deck.pptx")
            build_deck(manifest_path, pptx_path)
            assert os.path.exists(pptx_path), "PPTX file was not created"
            print(f"  Created: {pptx_path} ({os.path.getsize(pptx_path)} bytes)")

            # ---------------------------------------------------------------
            # Step 5: Run validate_deck -> verify passes
            # ---------------------------------------------------------------
            print("\n--- Step 5: Validate deck ---")
            render_dir = os.path.join(test_dir, "renders")
            report = validate_deck(pptx_path, output_dir=render_dir, json_output=True)
            print(f"  Passed: {report['passed']}")
            print(f"  Slide count: {report['slide_count']}")

            assert report["passed"], f"Validation failed: {report['issues']}"
            assert report["slide_count"] == len(manifest["slides"]), \
                f"Slide count mismatch: report={report['slide_count']} vs manifest={len(manifest['slides'])}"

            # ---------------------------------------------------------------
            # Step 6: Open .pptx and assert details
            # ---------------------------------------------------------------
            print("\n--- Step 6: Verify .pptx contents ---")
            prs = Presentation(pptx_path)

            # 6a: Slide count matches
            actual_count = len(prs.slides)
            expected_count = len(manifest["slides"])
            assert actual_count == expected_count, \
                f"Slide count: expected {expected_count}, got {actual_count}"
            print(f"  Slide count: {actual_count} (matches manifest)")

            # 6b: Background colour matches palette
            palette = PALETTES[palette_name]
            expected_bg = palette["bg"].upper()
            first_slide = prs.slides[0]
            bg_fill = first_slide.background.fill
            actual_bg = str(bg_fill.fore_color.rgb).upper()
            assert actual_bg == expected_bg, \
                f"Background colour mismatch: expected {expected_bg}, got {actual_bg}"
            print(f"  First slide bg: #{actual_bg} (matches palette '{palette_name}')")

            # 6c: At least one picture shape exists (image embedded)
            has_picture = False
            for slide in prs.slides:
                for shape in slide.shapes:
                    if isinstance(shape, Picture):
                        has_picture = True
                        break
                if has_picture:
                    break
            assert has_picture, "No embedded images found in deck"
            print(f"  Embedded images: found")

            print(f"\n{'='*60}")
            print("ALL ASSERTIONS PASSED")
            print(f"{'='*60}")

        except AssertionError as ae:
            print(f"\nASSERTION FAILED: {ae}")
            passed = False
        except Exception as ex:
            print(f"\nERROR: {ex}")
            passed = False
        finally:
            # Clean up test artifacts
            print(f"\n--- Cleanup ---")
            shutil.rmtree(test_dir, ignore_errors=True)
            print(f"  Removed: {test_dir}")

        return passed

    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


if __name__ == "__main__":
    success = run_e2e_test()
    sys.exit(0 if success else 1)
