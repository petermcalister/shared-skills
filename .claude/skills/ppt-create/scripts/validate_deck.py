#!/usr/bin/env python3
"""Validate a PowerPoint deck: structural checks + optional visual rendering.

Usage:
    cd .claude/skills/ppt-create/references && uv run python ../scripts/validate_deck.py --input <pptx> --output-dir <dir> --json
"""

import argparse
import inspect
import json
import logging
import os
import shutil
import subprocess
import sys
from typing import Any, Optional

from pptx import Presentation
from pptx.shapes.picture import Picture

logger = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format="%(levelname)s: %(message)s")


# ---------------------------------------------------------------------------
# Issue tracking
# ---------------------------------------------------------------------------

def _issue(slide_num: int, severity: str, message: str) -> dict[str, Any]:
    """Create a structured issue dict.

    Args:
        slide_num: 1-based slide number, or 0 for deck-level issues.
        severity: 'warning' or 'error'.
        message: Human-readable issue description.

    Returns:
        Dict with slide, severity, and message keys.
    """
    try:
        return {"slide": slide_num, "severity": severity, "message": message}
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# Structural validation
# ---------------------------------------------------------------------------

def validate_slide_count(prs: Presentation) -> list[dict[str, Any]]:
    """Check that slide count is reasonable (not 0, not >30).

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        List of issue dicts.
    """
    try:
        issues: list[dict[str, Any]] = []
        count = len(prs.slides)
        if count == 0:
            issues.append(_issue(0, "error", "Deck has no slides"))
        elif count > 30:
            issues.append(_issue(0, "warning", f"Deck has {count} slides (>30 may be too long)"))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def _shape_has_text(shape: Any) -> bool:
    """Return True if the shape contains non-empty text.

    Args:
        shape: A python-pptx shape object.

    Returns:
        True if the shape has text content.
    """
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            return len(text) > 0
        return False
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def _shape_is_picture(shape: Any) -> bool:
    """Return True if the shape is a picture/image.

    Args:
        shape: A python-pptx shape object.

    Returns:
        True if the shape is a Picture.
    """
    try:
        return isinstance(shape, Picture)
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def validate_empty_slides(prs: Presentation) -> list[dict[str, Any]]:
    """Check that every slide has at least one shape with text or an image.

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        List of issue dicts for empty slides.
    """
    try:
        issues: list[dict[str, Any]] = []
        for i, slide in enumerate(prs.slides, start=1):
            has_content = False
            for shape in slide.shapes:
                if _shape_has_text(shape) or _shape_is_picture(shape):
                    has_content = True
                    break
            if not has_content:
                issues.append(_issue(i, "warning", "Slide appears empty (no text or images)"))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def validate_images_embedded(prs: Presentation) -> list[dict[str, Any]]:
    """Check whether any slides contain embedded picture shapes.

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        List of issue dicts (informational if no images found).
    """
    try:
        issues: list[dict[str, Any]] = []
        has_any_image = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if _shape_is_picture(shape):
                    has_any_image = True
                    break
            if has_any_image:
                break
        if not has_any_image:
            issues.append(_issue(0, "warning", "No embedded images found in deck"))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def validate_title_slide(prs: Presentation) -> list[dict[str, Any]]:
    """Check that the first slide has large text (likely a title slide).

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        List of issue dicts.
    """
    try:
        issues: list[dict[str, Any]] = []
        if len(prs.slides) == 0:
            return issues
        first_slide = prs.slides[0]
        has_large_text = False
        for shape in first_slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.font.size and para.font.size.pt >= 30:
                        has_large_text = True
                        break
            if has_large_text:
                break
        if not has_large_text:
            issues.append(_issue(1, "warning", "First slide may not be a title slide (no large text found)"))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def validate_closing_slide(prs: Presentation) -> list[dict[str, Any]]:
    """Check that the last slide exists (closing slide present).

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        List of issue dicts.
    """
    try:
        issues: list[dict[str, Any]] = []
        if len(prs.slides) < 2:
            issues.append(_issue(0, "warning", "Deck has fewer than 2 slides — no distinct closing slide"))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def validate_text_lengths(prs: Presentation) -> list[dict[str, Any]]:
    """Warn if any text block exceeds 500 characters.

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        List of issue dicts for overly long text.
    """
    try:
        issues: list[dict[str, Any]] = []
        for i, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if len(text) > 500:
                        issues.append(_issue(i, "warning",
                                             f"Text block has {len(text)} chars (>500 may be too much for a slide)"))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def validate_bullet_count(prs: Presentation) -> list[dict[str, Any]]:
    """Warn if any slide has more than 5 bullet paragraphs in a single text frame.

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        List of issue dicts.
    """
    try:
        issues: list[dict[str, Any]] = []
        for i, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    paras = [p for p in shape.text_frame.paragraphs if p.text.strip()]
                    if len(paras) > 5:
                        issues.append(_issue(i, "warning",
                                             f"Text frame has {len(paras)} non-empty paragraphs (>5 bullets may be too dense)"))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def run_structural_validation(prs: Presentation) -> list[dict[str, Any]]:
    """Run all structural validation checks.

    Args:
        prs: An open python-pptx Presentation object.

    Returns:
        Combined list of all issue dicts.
    """
    try:
        issues: list[dict[str, Any]] = []
        issues.extend(validate_slide_count(prs))
        issues.extend(validate_empty_slides(prs))
        issues.extend(validate_images_embedded(prs))
        issues.extend(validate_title_slide(prs))
        issues.extend(validate_closing_slide(prs))
        issues.extend(validate_text_lengths(prs))
        issues.extend(validate_bullet_count(prs))
        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# Visual validation
# ---------------------------------------------------------------------------

def find_soffice() -> Optional[str]:
    """Find the soffice/LibreOffice executable on PATH or common locations.

    Returns:
        Path to soffice executable, or None if not found.
    """
    try:
        path = shutil.which("soffice")
        if path:
            return path
        # Common Windows locations
        for candidate in [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]:
            if os.path.exists(candidate):
                return candidate
        return None
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


def run_visual_validation(input_path: str, output_dir: str) -> list[dict[str, Any]]:
    """Convert .pptx to PNGs via LibreOffice and report dimensions.

    If LibreOffice is not installed, logs a warning and returns an empty list.

    Args:
        input_path: Absolute path to the .pptx file.
        output_dir: Directory to write rendered PNG files.

    Returns:
        List of issue dicts (informational).
    """
    try:
        issues: list[dict[str, Any]] = []
        soffice = find_soffice()
        if soffice is None:
            logger.warning("LibreOffice/soffice not found — skipping visual validation")
            issues.append(_issue(0, "warning", "Visual validation skipped: soffice not found"))
            return issues

        os.makedirs(output_dir, exist_ok=True)
        cmd = [soffice, "--headless", "--convert-to", "png", "--outdir", output_dir, input_path]
        result = subprocess.run(cmd, capture_output=True, text=True, shell=True, timeout=120)

        if result.returncode != 0:
            issues.append(_issue(0, "warning", f"soffice conversion failed: {result.stderr.strip()}"))
            return issues

        # Report dimensions of rendered PNGs
        try:
            from PIL import Image
            png_files = sorted([f for f in os.listdir(output_dir) if f.lower().endswith(".png")])
            for png_file in png_files:
                png_path = os.path.join(output_dir, png_file)
                with Image.open(png_path) as img:
                    w, h = img.size
                    logger.info("Rendered %s: %dx%d", png_file, w, h)
        except ImportError:
            logger.warning("Pillow not installed — cannot report PNG dimensions")
            # Still list the files
            png_files = sorted([f for f in os.listdir(output_dir) if f.lower().endswith(".png")])
            for png_file in png_files:
                logger.info("Rendered %s (dimensions unknown — install Pillow)", png_file)

        return issues
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# Main validation entry point
# ---------------------------------------------------------------------------

def validate_deck(input_path: str, output_dir: Optional[str] = None,
                  json_output: bool = False) -> dict[str, Any]:
    """Validate a PowerPoint deck and return a structured report.

    Args:
        input_path: Path to the .pptx file.
        output_dir: Optional directory for rendered PNGs.
        json_output: Whether to include JSON-formatted output.

    Returns:
        Report dict with file, slide_count, issues, and passed keys.
    """
    try:
        input_path = os.path.abspath(input_path)
        if not os.path.exists(input_path):
            raise FileNotFoundError(f"Input file not found: {input_path}")

        prs = Presentation(input_path)
        slide_count = len(prs.slides)

        # Structural validation
        issues = run_structural_validation(prs)

        # Visual validation (if output_dir provided)
        if output_dir:
            output_dir = os.path.abspath(output_dir)
            visual_issues = run_visual_validation(input_path, output_dir)
            issues.extend(visual_issues)

        has_errors = any(iss["severity"] == "error" for iss in issues)
        report: dict[str, Any] = {
            "file": os.path.basename(input_path),
            "slide_count": slide_count,
            "issues": issues,
            "passed": not has_errors,
        }

        # Print human-readable summary
        print(f"\n=== Deck Validation: {report['file']} ===")
        print(f"Slides: {slide_count}")
        if issues:
            for iss in issues:
                slide_ref = f"Slide {iss['slide']}" if iss["slide"] > 0 else "Deck"
                print(f"  [{iss['severity'].upper()}] {slide_ref}: {iss['message']}")
        else:
            print("  No issues found.")
        print(f"Result: {'PASSED' if report['passed'] else 'FAILED'}")

        # JSON output
        if json_output:
            print("\n--- JSON Report ---")
            print(json.dumps(report, indent=2))

        return report
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> None:
    """CLI entry point for validate_deck.py."""
    try:
        parser = argparse.ArgumentParser(description="Validate a PowerPoint deck.")
        parser.add_argument("--input", required=True, help="Path to input .pptx file")
        parser.add_argument("--output-dir", default=None, help="Directory for rendered PNGs")
        parser.add_argument("--json", action="store_true", dest="json_output",
                            help="Output JSON report to stdout")
        args = parser.parse_args()

        validate_deck(args.input, args.output_dir, args.json_output)
    except Exception as e:
        current_function = inspect.currentframe().f_code.co_name
        print(f"An error occurred in {current_function}: {e}")
        logger.warning(f"An error occurred in {current_function}: {e}")
        raise e


if __name__ == "__main__":
    main()
